#Import necessary libraries
import streamlit as st
from dotenv import load_dotenv
import os
from datetime import datetime,timedelta
import time as t
import openai
import ast
import pandas as pd
import base64
import re
import io
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from PIL import Image as PILImage
from streamlit_modal import Modal

#Send notification alert to phone
from twilio.rest import Client
import keys

#Library to generate .PDF file
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_JUSTIFY
from reportlab.lib.styles import ParagraphStyle
from reportlab.platypus import PageBreak, Image

#Import variables from custom modules
from explore_data import revenue_per_date, most_ordered_product, total_revenue, avg_wait_and_ready_time, order_per_date, merged_df, total_orders
from visualisation import sales_trend, best_product, orders_trend
from urllib.parse import unquote
from real_time import start_streaming_thread, start_inv_sim, seq, shared_state, start_summary_thread


# --- SESSION STATE INIT ---
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
if "username" not in st.session_state:
    st.session_state.username = ""
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "bot_typing" not in st.session_state:
    st.session_state.bot_typing = False

# --- CONFIG ---
st.set_page_config(page_title="Grab Login", page_icon="🟢", layout="wide")

# --- BACKGROUND IMAGE FUNCTION ---
def get_base64_image(image_path):
    with open(image_path, "rb") as img_file:
        encoded = base64.b64encode(img_file.read()).decode()
    return encoded

bg_image = get_base64_image("background.jpg")  # Use your image here

# Custom CSS for login and main page background image
st.markdown(f"""
    <style>
    .stApp {{
        background-image: url("data:image/jpeg;base64,{bg_image}");
        background-size: cover;
        background-position: center;
        background-repeat: no-repeat;
    }}

    div[data-testid="stForm"] {{
        background-color: rgba(0, 100, 0, 0.85);
        padding: 2rem;
        border-radius: 12px;
        color: white;
        box-shadow: 0 4px 12px rgba(0,0,0,0.3);
    }}

    label, .stTextInput label, .stPasswordInput label {{
        color: white !important;
    }}

    input {{
        border: 1px solid #00B14F !important;
        border-radius: 5px !important;
        color: black !important;
    }}

    div[data-testid="stForm"] button {{
        background-color: white !important;
        color: black !important;
        font-weight: bold;
        border: none;
        border-radius: 8px;
        padding: 0.5em 2em;
        box-shadow: 0px 2px 6px rgba(0, 0, 0, 0.2);
    }}

    div[data-testid="stForm"] button:hover {{
        background-color: #f0f0f0 !important;
        transition: 0.3s ease-in-out;
    }}

    div[data-testid="stPasswordInput"] > div {{
        position: relative;
    }}
    div[data-testid="stPasswordInput"] svg {{
        position: absolute;
        right: 10px;
        top: 38px;
        z-index: 10;
        color: #ffffff88;
    }}
    </style>
""", unsafe_allow_html=True)


#Function to generate powerpoint .ppt format slide
def save_chat_pptx():
    """Save chat history to text file with Q/A pairs and graph context"""
    if not st.session_state.chat_history:
        return None
    
    buffer = io.StringIO()
    pair_number = 1
    graph_counter = 0  # Separate counter for graph numbering
    
    # Skip initial bot message if exists
    start_index = 0
    if st.session_state.chat_history[0][0] == "Bot":
        start_index = 1
    
    for i in range(start_index, len(st.session_state.chat_history), 2):
        if i+1 >= len(st.session_state.chat_history):
            break  # Skip incomplete pairs
        
        question_entry = st.session_state.chat_history[i]
        answer_entry = st.session_state.chat_history[i+1]
        
        # Determine graph context
        if answer_entry[2] is not None:  # Check if graph exists
            graph_note = f"Graph{graph_counter}.png"
            graph_counter += 1
        else:
            graph_note = "None"
        
        buffer.write(f"{pair_number}. Question:\n{question_entry[1]}\n\n")
        buffer.write(f"Answer:\n{answer_entry[1]}\n")
        buffer.write(f"Graph: {graph_note}\n")
        buffer.write("-" * 50 + "\n\n")
        pair_number += 1
    
    buffer.seek(0)
    
    #Prompt to give instructions to OpenAi API for generating text on powerpoint slide
    prompt =f"""

    Context = {buffer.getvalue()}
        
    I need to make a presentation slide about the context i provided to you, write to me the structure of my powerpoint slides,
    for each page write exactly what text should be written on that specific slide,
    if the page only contain too little words, you can add some explanations if possible but not too long and dont hallucinate
    For the graphs, you can put it wherever you think it is necessary
    the structure of your output should be like this example,
    
    

    slide 1 :
    Pandas are animals that are endangered and we must protect them
    Same with Malayan Tigers
    --------------------------------------
    slide 2 :
    We must protect endangered animals
    animals are a valueable asset to our planet
    --------------------------------------
    slide 3 :
    .
    .
    .

    slide 1 will be the front page of the powerpoint therefore the content of the first slide should ALWAYS BE "Chat History Report"
    Ensure that each slide dont contain too little words
    YOUR OUTPUT SHOULD JUST BE THE STRUCTURE OF THE SLIDES AND NOTHING ELSE AND YOUR OUTPUT MUST STRICTLY FOLLOW THE FORMAT THAT I GAVE TO YOU
    """
    
    #Prompt to give instructions to OpenAi API for generating visuals on powerpoint slide
    structure_slide = get_ai_reply(prompt)
    prompt2 = f'''
{structure_slide}
this is the structure of a powerpoint slide i want you to make, i have already written out what to write on each page, I want you to write me a python code to make the powerpoint slide, 
for the image just read it from the current directory like Graph0.png, and also make sure that the graph is big enough but not too big untill it clashes with the words,
 make sure the slides and words are not out of the slides even partially and very beautiful and you can use themes, pick one from here

for the first slide which displays "Chat History Report", dont use the built-in title placeholder to display the "Chat History Report" title, instead, Manually add a textbox for the title  and put it at the center of the screen

theme_templates = [
        'theme_folder\Facet.pptx',          # Default
        'theme_folder\Gallery.pptx',                # Modern blue
        'theme_folder\Madison.pptx',             # Professional dark
        'theme_folder\Wisp.pptx',                  # Purple accents
        'theme_folder\Parallax.pptx',              # Green tones
    ], so pick one randomly

If the slide only have words and no image, make the font size between 18 to 24 depending on how packed the slide is (make sure that the words does not overflow outside of the slide)
If a slide have a graph in it make sure the fonts are in size of 20 or below depending on how packed the slide is


make sure that the words and the graphs are at the center of the slide and the words or graphs are not too big untill it overflows outside the slide 
The first slide is always the title of the entire slide so make it look like a header
to make this easier, you can hardcode it 
for each slide, delete the "Click to add title" box
store the whole powerpoint slide in a variable called prs
MAKE SURE THAT YOUR RESPONSE IS JUST PYTHON CODE AND NOTHING ELSE

at the end of the code, add this code

from pptx.enum.shapes import PP_PLACEHOLDER
for slide in prs.slides:
    # Make a list so we can safely remove while iterating
    for shape in list(slide.shapes):
        if (
            shape.is_placeholder 
            and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE
        ):
            # Remove the title placeholder from the slide XML
            slide.shapes._spTree.remove(shape._element)

to remove the text box "Click To Add Title"  for each slide

and also this code

# Access the underlying slide ID list
slide_id_list = prs.slides._sldIdLst  

# Remove the first entry (index 0)
first_slide_id = slide_id_list[0]
slide_id_list.remove(first_slide_id)

to remove the first slide because the first slide is for the template
DONT AND PLEASE NOT use MSO_ANCHOR cuz it will error
import necessary libraries first 

MAKE SURE THAT YOUR RESPONSE IS JUST PYTHON CODE AND NOTHING ELSE
    '''

    code = get_ai_reply(prompt2)
    code = re.sub(r'```python', '', code)
    code = re.sub(r'```', '', code)

    print("PPTX Code : ")

    with open("ppt_code.txt", "w") as f:
            f.write(code)

    exec_globals={}
    if is_code_safe(code):
        exec(code,exec_globals)
        prs = exec_globals['prs']
    else:
        print("There is something wrong")

    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer


#Appear a pop-up message to user
def new_day_modal():
    modal = Modal("Issue: Long wait times for drivers before pickup yesterday", key="new_day_modal")
    modal.open()

    if modal.is_open():
        with modal.container():
            st.write("""
                    Issue:

                    When: April 16th, 6:00 PM - 8:00 PM
                    
                    Problem: Drivers waited an average of 12 minutes after arriving, 5 minutes longer than usual.
                    
                    Effect: This caused an overall delivery delay, increasing the average time to 30 minutes, 10 minutes longer than expected.

                    Solution: Coordinate with kitchen or restaurant staff to ensure food is ready on time.""")
            if st.button("Coordinate with staff"):
                modal.close()


# --- Load valid users ---
df = pd.read_csv("crude_csv/merchant.csv")
valid_usernames = df["merchant_name"].tolist()

# --- LOGIN PAGE ---
if not st.session_state.logged_in:
    image = PILImage.open("logograbhijau.png")
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.image(image)

    st.markdown("<h1 style='text-align: center; color: #00B14F;'>Welcome to Grab Merchant Portal</h1>", unsafe_allow_html=True)

    with st.form("login_form"):
        st.markdown("### Please login to continue")
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")
        login_button = st.form_submit_button("Login")

        if login_button:
            if username in valid_usernames and password == "123":
                st.session_state.logged_in = True
                st.session_state.username = username
                st.success(f"✅ Welcome, {username}!")
                st.query_params.shop = username
                st.session_state.show_modal = True
                st.rerun()
            elif username not in valid_usernames:
                st.error("❌ Invalid username.")
            else:
                st.error("❌ Invalid password.")
            
else:

    # Modal after login — shows once
    if st.session_state.show_modal:
        modal = Modal("🎉 Welcome!", key="welcome_modal")
        with modal.container():
            st.markdown("""
                Issue:

                        When: April 16th, 6:00 PM - 8:00 PM
                        Problem: Drivers waited an average of 12 minutes after arriving, 5 minutes longer than usual.
                        Effect: This caused an overall delivery delay, increasing the average time to 30 minutes, 10 minutes longer than expected.

                        Solution:

                        Coordinate with kitchen or restaurant staff to ensure food is ready on time.
            """, unsafe_allow_html=True)
            if st.button('Coordinate with staff'):
                st.markdown('Thank you for taking action.')
                t.sleep(2.5)
            st.session_state.show_modal = False


    #Function to generate PDF .pdf format slide
    def generate_chat_pdf():
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter,
                            rightMargin=72, leftMargin=72,
                            topMargin=72, bottomMargin=18)
        
        styles = getSampleStyleSheet()
        styles.add(ParagraphStyle(name='Justify', alignment=TA_JUSTIFY))
        Story = []
        
        # Add title
        title_style = styles["Title"]
        title = Paragraph("Chat History Report", title_style)
        Story.append(title)
        Story.append(Spacer(1, 24))
        
        qna_count = 0  # Track Q&A pairs per page
        img_counter = 0  # Track graph images
        
        chat_history = st.session_state.chat_history
        
        # Skip initial bot message if exists
        start_index = 0
        if chat_history and chat_history[0][0] == "Bot":
            start_index = 1
        
        for i in range(start_index, len(chat_history), 2):
            if i+1 >= len(chat_history):
                break
            
            if qna_count == 2:
                # Add page break
                Story.append(Spacer(1, 12))
                Story.append(PageBreak())
                qna_count = 0
            
            user_entry = chat_history[i]
            bot_entry = chat_history[i+1]
            
            # Question
            question_text = f"<b>Question:</b> {user_entry[1]}"
            p_question = Paragraph(question_text, styles["BodyText"])
            Story.append(p_question)
            Story.append(Spacer(1, 12))
            
            # Answer
            answer_text = f"<b>Answer:</b> {bot_entry[1]}"
            p_answer = Paragraph(answer_text, styles["BodyText"])
            Story.append(p_answer)
            Story.append(Spacer(1, 24))
            
            # Handle graph
            if bot_entry[2] is not None:
                try:
                    img_path = f"graph{img_counter}.png"
                    img = Image(img_path, width=5*inch, height=3*inch)
                    Story.append(img)
                    Story.append(Spacer(1, 24))
                    img_counter += 1
                    
                    # If image takes too much space, force new page
                    if qna_count == 1:
                        Story.append(PageBreak())
                        qna_count = 0
                except Exception as e:
                    error_text = Paragraph("<i>[Graph unavailable]</i>", styles["Italic"])
                    Story.append(error_text)
            
            qna_count += 1
        
        doc.build(Story)
        buffer.seek(0)
        return buffer.getvalue()

    if "initial_model" not in st.session_state:
        st.session_state.initial_model = False


    FORBIDDEN_NAMES = { "sys", "subprocess", "shutil", "__import__"}

    #Function to determine the executed code is save
    def is_code_safe(code: str) -> bool:
        try:
            tree = ast.parse(code)
        except SyntaxError as e:
            print("Syntax error in code:", e)
            return False

        for node in ast.walk(tree):
            # Check for dangerous imports
            if isinstance(node, (ast.Import, ast.ImportFrom)):
                names = [alias.name for alias in node.names]
                if any(name.split('.')[0] in FORBIDDEN_NAMES for name in names):
                    print(f"Found forbidden import in code: {names}")
                    return False

            # Check for dangerous function calls
            if isinstance(node, ast.Call):
                if isinstance(node.func, ast.Name) and node.func.id in {"eval", "exec", "__import__"}:
                    print(f"Found dangerous function call: {node.func.id}")
                    return False

        return True

    load_dotenv()

    os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"

    # Function to get AI response using OpenAI
    openai.api_key = os.environ.get("OPENAI_API_KEY")

    #Function to call OpenAI API
    def get_ai_reply(query):
        response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": query}
            ],
            temperature=0.7
        )

        return response.choices[0].message.content

    start_streaming_thread('Bagel Bros')
    start_streaming_thread('Noodle Nest')

    if "inv_status" not in st.session_state:
        st.session_state.inv_status = None

    start_inv_sim(seq)

    new_status = shared_state.get("inv_status")
    if new_status is not None:
        st.session_state.inv_status = new_status

    if st.button('Click to Refresh'):
        st.rerun()

    # Sidebar layout for the chatbot
    shop_param = st.query_params.get_all('shop')
    if shop_param:
        shop = unquote(shop_param[0])
    else:
        shop = 'Bagel Bros'

    time_param = st.query_params.get_all('time')
    if time_param:
        time = unquote(time_param[0])
    else:
        time = 'Today'

    if "shop" not in st.session_state:
        st.session_state.shop = 'Bagel Bros'
    else:
        shop = st.session_state.shop

    if "time" not in st.session_state:
        st.session_state.time = 'Today'
    else:
        time = st.session_state.time

    df_merged_live = merged_df(shop, None, live=True)
    df_past_data = merged_df(shop, None, all=True)
    df_merged = pd.concat([df_past_data, df_merged_live], ignore_index=True)

    additional_context = df_merged["item_name"].unique()

    datetime_columns = ['order_time', 'driver_arrival_time', 'driver_pickup_time', 'delivery_time']

    for col in datetime_columns:
        df_merged[col] = pd.to_datetime(df_merged[col],errors='coerce')

    #Give the history context to the chatbot
    def format_chat_history(chat_history):
        """Formats chat history into a string of Q/A pairs"""
        history_str = ""
        pair_num = 1
        
        # Skip initial bot message if exists
        start_index = 0
        if chat_history[0][0] == "Bot":
            start_index = 1
        
        for i in range(start_index, len(chat_history), 2):
            if i+1 >= len(chat_history):
                break
            
            question = chat_history[i][1]
            answer = chat_history[i+1][1]
            
            history_str += f"Conversation History {pair_num}:\n"
            history_str += f"Question: {question}\n"
            history_str += f"Answer: {answer}\n\n"
            pair_num += 1
        
        return history_str

    #Give the history context to the chatbot
    def is_followup(query):
        followup_cues = ["what about", "how about", "and", "also", "other", "more", "it", "them", "those"]
        return any(cue in query.lower() for cue in followup_cues)

    #Prompt OpenAI API to response within the given context
    def chat_bot_output(user_input):
        chat_history = format_chat_history(st.session_state.chat_history)
        summary_chat = st.session_state.memory_summary
        #user_input = "Based on the data for the past 14 days (today 17th June 2023), which hour usually causes operational bottlenecks and how to avoid that efficiently"
        prompt = f'''
        This merchant name is {shop}.

        You are a RAG LLM Agent that will answer questions about a csv file, the csv file is already readed and right now the csv file is in a variable called df_merged . This dataset is data about food delivery services
        The dataset contains columns which is :

        - "order_id" (this is the id of each order, this is not a primary key because here some rows have the same order_id but different item_id which means that the order consist of more than one meal)
        - "order_time" (this is a datetime column that shows when the order is placed by the customer)
        - "driver_arrival_time" (this is a datetime column that shows when the delivery guy arrives at a restaurant)
        - "driver_pickup_time" (this is a datetime column that shows when the delivery guy able to pick up the food from the restaurant to send to the customer)
        - "delivery_time" (this is a datetime column that shows when the delivery guy arrives at the customer's house to deliver the food)
        - "eater_id" (this is the id of each customer)
        - "order_ready" (this just the period (in minutes) between the column driver_pickup_time and driver_order_time)
        - "driver_arrival_period" (this is the period (in minutes) between the column order_time and driver_arrival_time)
        - "driver_waiting_time" (this is the period (in minutes) between the column driver_arrival_time and driver_pickup_time)
        - "item_id" (this is the id of each meal/items)
        - merchant_id_x (this is the id of the merchant/restaurant)
        - cuisine_tag  (this is the category of each meal, like "western","breakfast",etc..)
        - item_name (this is the name of the meal/item)
        - item_price (this is the price of the meal/item)
        - ingredients (this is an array of ingredients used to make the meal but this is a string column, example values are like "["Onions","Rice","Egg"]" so you need to handle that)

        for your information, the merchant_id_x for Bagel Bros is 3e2b6, and the city_id for Noodle Nest is 9b5a0.

        we also have a csv file called merchant.csv which stores merchants information (ONLY USE THIS DATA WHEN THE USER IS ASKING ABOUT OTHER MERCHANTS) and the file is located in a folder called processed_csv, so to read it just do pd.read_csv("processed_csv\dated_merchant.csv")
        to answer questions about competitors, use this dataset which contains columns such as:
        - merchant_id (this is the id of each merchants)
        - merchant_name (which is the merchant name)
        - join_date (which is the date when the merchant is established or registered, formatted as dd-mm-yyyy)
        - city_id (which is an integer column where each integer represents a city the merchant is in)

        When generating Python code to answer queries, make sure to parse the join_date column as datetime using pd.to_datetime(..., format='%d-%m-%Y') so that date-based filtering works correctly.

        for your information, the city_id for Bagel Bros is 8, and the city_id for Noodle Nest is 5.

        city_id in dated_merchant.csv IS NOT THE SAME AS merchant_id_x from df_merged. city_id is the id of the city the merchant is located in, WHILE merchant_id_x is the merchant identifier that identifies the merchant.

        the city Bagel Bros located in has id of 8 (city_id = 8), but Bagel Bros merchant_id IS NOT 8. DO NOT USE THEM INTERCHANGABLY. Same goes with Noodle Nest.

        additional information, These are all the products that they sold (unique values in the item_name column ):{additional_context}

        Based on this, write me a python code to answer this question {user_input} , at the end of the code store the final output in a variable called Final_Output (If you think that the question cannot be answered by the data, just say "Sorry I dont know") ,
        This  Final_Output variable will be send to an api to OpenAI to do some analysis so it cannot be None so make sure th structure of this variable is suitable for analysis
        This Final_Output also CANNOT BE A SENTENCE because we cant perform data analysis using a sentence,
        display your final output using print statements, or display a matplotlib graph (if the user wants a graph)
        If you displayed a graph, put it in a variable called Final_Graph (in plt format), if you are not displaying a graph, just set Final_Graph = None
        Ensure that the visual elements produced by the code is able to effectively enhance understanding of the data, even by non-technical people.
        If the user asks to do predictions and ML models are involved, include a gridsearch cv code to ensure 
        Import necessary Libraries and Your Output should only be PYTHON CODE and nothing else
        For your information this dataset only contains data in the year 2023, and today is 17th June 2023
        
        This is the recent chat history between you and the user, USE ONLY WHEN NECESSARY TO ANSWER THE USER QUESTION:
        {chat_history}
        This is the summarized older chat history between you and the user, USE ONLY WHEN NECESSARY TO ANSWER THE USER QUESTION:
        {summary_chat}
        If the current user input appears to be a follow-up (e.g. uses words like "what about", "other", "those", "more", "it"), then refer back to the previous chat message and the Final_Output from the previous query. Assume the user wants to build on the last insight.
        '''

        if is_followup(user_input):
            previous_output = st.session_state.get("previous_final_output", "")
            prompt += f"\nFor context, the previous output from a similar query was:\n{previous_output}\n"

        extra_string =f'''

import openai
import os
import matplotlib.pyplot as plt

openai.api_key = os.environ.get("OPENAI_API_KEY")

# Case 1: If Final_Output is a matplotlib figure (assumed to be shown already)
if isinstance(Final_Output, plt.Figure):
    print("Case 1")
    image_path = "plot_output.png"
    Final_Output.savefig(image_path)
    final_output_str = "A graph has been generated and saved as 'plot_output.png'."
    import base64
    client = OpenAI(
        api_key=os.getenv("OPENAI_API_KEY")
    )

    prompt = "Based on this image, analyse it and give a proper answer to the question '{user_input}'. For your information, my merchant name is {shop}. (answer the question using the language that is used in the query) do not output python code bcuz merchants cant understand"

    # Function to encode the image
    def encode_image(image_path):
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode("utf-8")
        
    # Encode the image
    base64_image = encode_image(image_path)

    # Create chat completion
    try:
        chat_completion = client.chat.completions.create(
            model="gpt-4o",
            max_tokens=300,
            messages=[
                {{
                    "role": "user",
                    "content": [
                        {{
                            "type": "text",
                            "text": prompt
                        }},
                        {{
                            "type": "image_url",
                            "image_url": {{
                                "url": f"data:image/jpeg;base64,{{base64_image}}"
                            }}
                        }}
                    ]
                }}
            ]
        )
        print("----------------------------------")
        Final_Reply = chat_completion.choices[0].message.content
    except Exception as e:
        print(f"An error occured: {{e}}")
    
# Case 2: If Final_Output is a plot shown via plt (stateful)


# Case 3: If it's just textual or tabular data
else:
    print("Case 3")
    final_output_str = str(Final_Output)
    query = f"""Based on final output : {{Final_Output}}, provide a proper answer to the query which is '{user_input}'. 
    For your information, my merchant name is {shop}. (answer the question using the language that is used in the query), do not output python code bcuz merchnts cant understand, if the question is not related at all about sales, just answer i dont know"""
    print(query)
    print("----------------------------------")
    Final_Reply = get_ai_reply(query)
    '''

        Final_Output = "An error occured"
        Final_Graph = None


        #user_input = "Produk apa paling laku do"
        code = get_ai_reply(prompt)
        code = code + extra_string

        code = re.sub(r'```python', '', code)
        code = re.sub(r'```', '', code)
        with open("logs.txt", "w",encoding = "utf-8") as f:
                f.write(code)


        # Remove code block markers
        code = re.sub(r'```python', '', code)
        code = re.sub(r'```', '', code)

        with open("logs.txt", "w",encoding = "utf-8") as f:
            f.write(code)


        exec_globals={}
        try:
            if is_code_safe(code):
                exec_globals["df_merged"] = df_merged
                exec_globals["get_ai_reply"] = get_ai_reply
                exec(code,exec_globals)
                Final_Output = exec_globals.get('Final_Reply')
                Final_Graph = exec_globals.get('Final_Graph')
            else:
                print("There is something wrong")
        except Exception as e:
            print(f"Error: {e}")


        print (Final_Output,Final_Graph)

        return Final_Output,Final_Graph

    # Session state initialization
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    if "bot_typing" not in st.session_state:
        st.session_state.bot_typing = False
    if "graph_i" not in st.session_state:
        st.session_state.graph_i = 0
    if "chat_hist_summary_pool" not in st.session_state:
        st.session_state.chat_hist_summary_pool = []
    if "memory_summary" not in st.session_state:
        st.session_state.memory_summary = None

    # Ensure bot gives an opening message
    if not st.session_state.chat_history:
        timestamp = datetime.now().strftime("%H:%M:%S")
        st.session_state.chat_history.append(("Bot", "Hello! How can I assist you today?",None, timestamp))

    # Simulate handling user messages
    def handle_user_message():
        user_input = st.session_state.chat_input.strip()
        
        if user_input:
            timestamp = datetime.now().strftime("%H:%M:%S")
            st.session_state.chat_history.append(("You", user_input,None, timestamp))
            st.session_state.chat_hist_summary_pool.append(("You", user_input,None, timestamp))

            st.session_state.chat_input = ""
            st.session_state.bot_typing = True
            response,graph = chat_bot_output(user_input)

            st.session_state.chat_history.append(("Bot", response,graph, datetime.now().strftime("%H:%M:%S")))
            st.session_state.chat_hist_summary_pool.append(("Bot", response,graph, datetime.now().strftime("%H:%M:%S")))

            if graph is not None :
                graph.savefig(f"graph{st.session_state.graph_i}.png", dpi=300, bbox_inches="tight")
                st.session_state.graph_i=st.session_state.graph_i+1

            st.session_state.bot_typing = False

    if len(st.session_state.chat_hist_summary_pool) > 1:
        start_summary_thread()
        print(len(st.session_state.chat_hist_summary_pool))
    else:
        print('summary not yet brada')

    # Custom CSS for the app layout
    st.markdown("""
        <style>
            body, .stApp {
                background-color: #fffff;
                color: black;
            }
            [data-testid="stSidebar"] {
                background-color: #dbf0ed;
                color: #08543c;
                border-radius: 15px;
            }
            [data-testid="stTextInput"] div div input {
                        background-color: white;
                        border: 2px solid white;
                        color: black;
                    }
                
            /* Main button styling */
            div[data-testid="stButton"] > button {
                background-color: #10543c !important;
                border: 2px solid white !important;
                color: white !important;
                transition: all 0.3s ease !important;
            }

            /* Hover effects */
            div[data-testid="stButton"] > button:hover {
                background-color: #0d4530 !important;
                border: 2px solid white !important;
                color: white !important;
                filter: brightness(90%);
            }

            /* Active/click state */
            div[data-testid="stButton"] > button:active {
                background-color: #0a3726 !important;
            }

            /* Focus state */
            div[data-testid="stButton"] > button:focus {
                box-shadow: 0 0 0 0.2rem rgba(255, 255, 255, 0.5) !important;
            }

            /* Disabled state */
            div[data-testid="stButton"] > button:disabled {
                background-color: #10543c88 !important;
                border-color: #ffffff88 !important;
                }
            [data-testid="stTextInput"] div div input {
                        background-color: white;
                        border: 2px solid white;
                        color: black;
            }
            .chat-box {
                max-height: 250px;
                overflow-y: auto;
                padding: 10px;
                background-color: #08543c;  /* new background color */
                border-radius: 10px;
                margin-bottom: 10px;
                color: white;  /* ensures text inside is white */
                border: 2px solid white;
            }
            .message {
                padding: 8px 12px;
                margin: 6px 0;
                border-radius: 10px;
                max-width: 90%;
                word-wrap: break-word;
                
            }
            .user {
                background-color: #0A9830;
                color: white;
                text-align: right;
                margin-left: auto;
                
                
            }
            .bot {
                background-color: #ffffff;
                color: black;
                text-align: left;
                margin-right: auto;
                
                
            }
            .timestamp {
                font-size: 10px;
                color: black;
            }
            
        </style>
    """, unsafe_allow_html=True)

    #Config mathplotlib graph to base64 image
    def fig_to_base64_img(fig):
        buf = io.BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight")
        buf.seek(0)
        encoded = base64.b64encode(buf.read()).decode("utf-8")
        return f"data:image/png;base64,{encoded}"

    #Allow chatbot to display on sidebar
    with st.sidebar:
        st.markdown("""
        <h1 style='text-align: center; font-size: 32px;'>🤖 Mex Assistant</h1>
        """, unsafe_allow_html=True)
        st.markdown(
        "<hr style='border: 1px solid white;'>",
        unsafe_allow_html=True
        )

        # Get query parameters from URL
        times = ["Today", "This Week", "This Month"]

        # Retrieve the shop query parameter from the URL (if it exists)
        if "time" not in st.query_params:
            # Default to "Bagel Bros" if no query parameter exists
            st.query_params.time = "Today"

        st.session_state.setdefault("times", times)

        #Changes on dropdown will effect the dashboard
        def update_options():
        # move the selected option to the front of the list if it is not already
            if st.session_state.selected_time != st.session_state.times[0]:
                st.session_state.times.remove(st.session_state.selected_time)
                st.session_state.times.insert(0, st.session_state.selected_time)

        st.markdown(
            """
            <style>
            /* Target the label above the selectbox */
            section div[data-testid="stSelectbox"] label {
                color: black !important;
                font-size: 16px !important;
                font-weight: bold !important;
            }
            
            div[data-baseweb="select"] > div {
                background-color: #D3D3D3;  /* Light blue */
                color: black;               /* Text color */
            }
            </style>
            """,
            unsafe_allow_html=True
        )

        time = st.selectbox(
            label='Select an option',
            options=st.session_state.times,
            key="selected_time",
            on_change=update_options,
        )

        # Update the shop value in query params
        st.session_state.time = time
        st.query_params.time = time

        today = "2023-06-17"

        #Set date
        def get_week_dates(date_str, fmt="%Y-%m-%d"):
            # Parse into a date
            dt = datetime.strptime(date_str, fmt).date()
            # In ISO, Monday=1 … Sunday=7; we want Sunday=0 shift
            days_to_sunday = dt.isoweekday() % 7
            # Find that week’s Sunday
            start_sun = dt - timedelta(days=days_to_sunday)
            # Build list from Sunday → Saturday
            return [(start_sun + timedelta(days=i)).strftime(fmt) for i in range(7)]

        order_goal = 100
        if(st.session_state.time=="Today"):
            print("yabadabadoo")
            print(st.query_params)
            if(st.query_params['shop']=="Bagel Bros"):
                print("yabadabadoo1")
                df= pd.read_csv("goal_bagel_fin.csv")
                row = df[df["Date"]==today]
                order_goal =  row.iloc[0]["Goal"]
            if(st.query_params['shop']== "Noodle Nest"):
                print("yabadabadoo2")
                df= pd.read_csv("goal_noodle.csv")
                row = df[df["Date"]==today]
                order_goal =  row.iloc[0]["Goal"]

        elif (st.session_state.time == "This Week"):
            print("yabadabadoo3")
            if(st.query_params['shop'] == "Bagel Bros"):
                df= pd.read_csv("goal_bagel_fin.csv")
                week_list = get_week_dates(today)
                week_rows = df[df["Date"].isin(week_list)]
                sum = week_rows["Goal"].sum()
                order_goal = sum
            if(st.query_params['shop'] =="Noodle Nest"):
                df= pd.read_csv("goal_noodle.csv")
                week_list = get_week_dates(today)
                week_rows = df[df["Date"].isin(week_list)]
                sum = week_rows["Goal"].sum()
                order_goal = sum
        else:
            if(st.query_params['shop']=="Bagel Bros"):
                df= pd.read_csv("goal_bagel_fin.csv")
                df2=df.copy()
                df2["Date"] = pd.to_datetime(df2["Date"])
                today_date = pd.to_datetime(today)
                current_month = today_date.month
                sum = df2[df2['Date'].dt.month == current_month]['Goal'].sum()
                order_goal = sum
            if(st.query_params['shop']=="Noodle Nest"):
                df= pd.read_csv("goal_noodle.csv")
                df2=df.copy()
                df2["Date"] = pd.to_datetime(df2["Date"])
                today_date = pd.to_datetime(today)
                current_month = today_date.month
                sum = df2[df2['Date'].dt.month == current_month]['Goal'].sum()
                order_goal = sum 

        # Scrollable chat history box
        with st.container():
            i=0
            #st.markdown('<div class="chat-box" style="max-height: 250px; overflow-y: auto;">', unsafe_allow_html=True)
            for sender, message, graph, timestamp in st.session_state.chat_history:
                sender_class = "user" if sender == "You" else "bot"
                align = "right" if sender == "You" else "left"
                if graph is not None:
                    col1, col2 = st.columns([6, 1])
                    with col1:
                        st.image(f"graph{i}.png",use_container_width=True)
                        i=i+1
                    with col2:
                        st.write("")

                st.markdown(
                    f"""
                    <div class="message {sender_class}">
                        <div>{message}</div>
                        <div class="timestamp" style="text-align: {align};">{timestamp}</div>
                    </div>
                    """,
                    unsafe_allow_html=True
                )
            st.markdown('</div>', unsafe_allow_html=True)

        # Text input for user query
        st.markdown(
        "<hr style='border: 1px solid white;'>",
        unsafe_allow_html=True
        )

        st.text_input("",placeholder="Ask the chatbot here", key="chat_input", value=st.session_state.get("chat_input", ""), on_change=handle_user_message)

        st.markdown(
                """
                <style>
                /* Targeting the download button inside its container */
                div.stDownloadButton button {
                    background-color: #10543c;
                    color: white;
                    border: 1px solid white;
                    padding: 0.75rem 1.5rem;  /* Adjust padding as needed */
                    font-size: 1rem;
                }
                div.stDownloadButton button:hover {
                    background-color: #0e4532;  /* Slightly darker on hover */
                }
                </style>
                """,
                unsafe_allow_html=True
            )
        col1, col2 = st.columns(2)
        with col1:
            if st.button("Download Chat As PDF", use_container_width=True):
                pdf_bytes = generate_chat_pdf()
                st.download_button(
                    label="Confirm Download",
                    data=pdf_bytes,
                    file_name="chat_history.pdf",
                    mime="application/pdf",
                    use_container_width=True
                )
        with col2:
            if st.button("Download Chat as PPT", use_container_width=True):
                pptx_file = save_chat_pptx()
                if pptx_file:
                    st.download_button(
                        label="Confirm PPTX Download",
                        data=pptx_file,
                        file_name="chat_history.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                else:
                    st.warning("Chat history is empty!")

    #CSS style for displaying header on the dashboard
    def bordered_card(content, bg_color="#ffffff", text_color="black"):
        st.markdown(
            f"""
            <div style="
                border: 2px solid #28a745;
                border-radius: 15px;
                padding: 20px;
                background-color: {bg_color};
                color: {text_color};
                height: 100%;
            ">
                {content}
            </div>
            """,
            unsafe_allow_html=True
        )

    # Reusable card layout
    def bordered_card_metric(title, value, delta=None, bg_color="#ffffff", text_color="black"):
        # Arrow logic for delta
        delta_display = ""
        if delta is not None:
            arrow = "▲" if str(delta).startswith("+") or not str(delta).startswith("-") else "▼"
            delta_color = "green" if arrow == "▲" else "red"
            delta_display = f"<small style='color:{delta_color}; font-size: 0.9em;'>{arrow} {delta}</small>"

        st.markdown(
            f"""
            <div style="
                border: 2px solid #28a745;
                border-radius: 15px;
                padding: 20px;
                background-color: {bg_color};
                color: {text_color};
                height: 100%;
                text-align: center;
            ">
                <p style='margin:0; font-weight:bold;'>{title}</p>
                <h3 style='margin:0;'>{value}</h3>
                {delta_display}
            </div>
            """,
            unsafe_allow_html=True
        )


    # Top bar
    top_left, spacer, profile, settings, logout = st.columns([3, 5, 1, 1, 1])
    with top_left:
        st.image("grab2.png", width=150)
    with profile:
        st.write("**Profile**")
    with settings:
        st.write("**Settings**")
    with logout:
        if st.button("Logout"):
            st.session_state.logged_in = False
            st.session_state.username = ""
            st.rerun()

    st.markdown("---")

    # Dashboard Title
    st.markdown('<h1 style="color:black;">Seller Dashboard</h1>', unsafe_allow_html=True)


    # Top row
    driver_wait, order_prep = avg_wait_and_ready_time(shop)


    progress_text = f"Order Goals For {time}: {order_goal}"
    orders1, orders2, orders3 = total_orders(shop, time=time, live=True)

    if time == 'Today':
        orders = orders1
    else:
        orders = orders2

    # Compute a 0.0–1.0 fraction and cap at 1.0
    percentage_progress = orders / order_goal
    if percentage_progress > 1:
        percentage_progress = 1.0

    # Convert to integer percent for display
    percent = int(percentage_progress * 100)

    st.markdown(f"""
    <div style="
        position: relative;
        background-color: #e0e0e0;
        border-radius: 8px;
        width: 100%;
        height: 24px;
        overflow: hidden;
    ">
    <!-- the fill -->
    <div style="
        width: {percent}%;
        background-color: #08543c;
        height: 100%;
    "></div>

    <!-- overlay text (centered) -->
    <div style="
        position: absolute;
        top: 0;
        left: 0;
        width: 100%;
        height: 100%;
        display: flex;
        align-items: center;
        justify-content: center;
        color: white;
        font-weight: bold;
        pointer-events: none;
    ">
        {progress_text} ({percent}%)
    </div>
    </div>

    <!-- bottom text (below the bar) -->
    <p style="
        text-align: center;
        margin-top: 6px;
        font-size: 0.9em;
    ">
                
    </p>
    """, unsafe_allow_html=True)

    if time == 'Today':
        time_comp = 'Yesterday'
    elif time == 'This Week':
        time_comp = 'Last Week'
    else:
        time_comp = 'Last Month'

    col1, col6, col7, col2 = st.columns(4)
    base_rev, revenue, diff_rev = total_revenue(shop, live=True, time=time_comp)
    base_rev = "{:,.2f}".format(base_rev)
    if diff_rev < 0:
        diff_rev = "-RM{:,.2f}".format(diff_rev)
    else:
        diff_rev = "RM{:,.2f}".format(diff_rev)

    with col1:
        bordered_card_metric(f"{time}'s Total Sale", f"RM{base_rev}", f"{diff_rev} Since {time_comp}")

    base_orders, orders, diff_orders = total_orders(shop, time=time_comp, live=True)
    with col6:
        bordered_card_metric(f"{time}'s Total Orders", f"{base_orders}", f"{diff_orders} Since {time_comp}")

    order_prep = "{:,.2f}".format(order_prep)
    with col7:
        bordered_card_metric("Time to Prepare Order", f"{order_prep} min")

    driver_wait = "{:,.2f}".format(driver_wait)
    with col2:
        bordered_card_metric("Driver Waiting Time", f"{driver_wait} min")
        
    if "noti" not in st.session_state:
        st.session_state.noti = []


    #NOTIFICATION ALERT 1
    if st.session_state.inv_status == 0:
        messageToUser1 = """
        🔄 Flour inventory is running out! Do restock ASAP
        
        Issue:

            Flour inventory is running out. Please restock ASAP

            Solution: Buy more supplies for your shop.
        
        """
        
        st.toast(messageToUser1)
        st.session_state.noti.append(messageToUser1)
        
        client = Client(keys.account_sid, keys.auth_token)

        message = client.messages.create(
        body = messageToUser1,
        from_ = keys.twilio_number,
        to = keys.my_phone_number
        )

        print(message.body)

        modal = Modal("🔄 Flour Restock Needed", key="restock_modal")
        with modal.container():
            st.markdown("""
                Issue:

                        Flour is running out. Please restock ASAP

                        Solution:

                        Buy more supplies for your shop.
            """, unsafe_allow_html=True)
            if st.button('Buy more supplies'):
                st.markdown('Thank you for taking action.')
                t.sleep(2.5)
            

    st.write(" ")
    st.write(" ")
    st.write(" ")

    col3, col4 = st.columns([3, 1])
    with col3:
        bordered_card("<strong>Orders Trend</strong><br>")
        order_data = order_per_date(shop, time)
        plt = orders_trend(order_data, time)
        st.pyplot(plt)
    with col4:
        bordered_card("<strong>Product Performance</strong><br>")
        quantity_data = most_ordered_product(shop, time)
        plt = best_product(quantity_data, time)
        st.pyplot(plt)

    # Middle row
    col8, col5 = st.columns([3, 1])
    with col8:
        bordered_card("<strong>Revenue Trend</strong><br>")
        tr_data = revenue_per_date(shop, time)
        plt = sales_trend(tr_data, time)
        st.pyplot(plt)

    with col5:
        noti_text = "No new notifications."

        if len(st.session_state.noti) > 0:
            noti_text = ""
            for i in range(len(st.session_state.noti)):
                noti_text = "<br>".join(st.session_state.noti)
                
        bordered_card(f"<strong>Notification Centre</strong><br><p>{noti_text}</p>", bg_color="#28a745", text_color="white")

    if st.session_state.username == 'Bagel Bros':
        csv = 'bagel_fin'
    else:
        csv = 'noodle'

    #Predict graph for next 5 days
    def predictions_graph(date_str, csv_path=f"goal_{csv}.csv"):
        # 1) Load and normalize dates
        df = pd.read_csv(csv_path)
        df["Date"] = pd.to_datetime(df["Date"]).dt.strftime("%Y-%m-%d")
        
        # 2) Locate the first matching row
        matches = df.index[df["Date"] == date_str].tolist()
        if not matches:
            raise ValueError(f"No data for date {date_str}")
        start_idx = matches[0]
        
        # 3) Slice out this row + next 6 rows (total 7)
        df_week = df.iloc[start_idx : start_idx + 7]
        
        # 4) Plot with custom bar color
        fig, ax = plt.subplots(figsize=(8, 4))
        ax.bar(
            df_week["Date"],
            df_week["Goal"],
            color="#08543c"           # <-- your custom green
        )
        ax.set_xlabel("Date")
        ax.set_ylabel("Predicted Order")
        plt.xticks(rotation=45, ha="right")
        plt.tight_layout()
        
        st.pyplot(plt)

    if(st.session_state.time=="Today"):
        bordered_card("<strong>Order Prediction For The Next 7 days</strong><br>")
        predictions_graph(today)